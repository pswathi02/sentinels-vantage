#!/usr/bin/env python3
"""
Build the Vantage / Team Sentinels submission deck as a .pptx that imports
cleanly into Google Slides (File ▸ Import slides, or just upload to Drive and
open with Google Slides).

Run:  python3 submission/build_deck.py
Out:  submission/Vantage-Sentinels-Submission.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ---- palette (matches the app's dark theme) ----
BG        = RGBColor(0x0B, 0x0E, 0x14)
BG_CARD   = RGBColor(0x16, 0x1C, 0x29)
BG_ELEV   = RGBColor(0x12, 0x17, 0x22)
BORDER    = RGBColor(0x2A, 0x35, 0x49)
TEXT      = RGBColor(0xE8, 0xED, 0xF6)
TEXT_DIM  = RGBColor(0x9F, 0xB0, 0xC8)
TEXT_FAINT= RGBColor(0x6B, 0x7A, 0x93)
ACCENT    = RGBColor(0x4F, 0x9D, 0xFF)
ACCENT2   = RGBColor(0x7C, 0x5C, 0xFF)
GREEN     = RGBColor(0x36, 0xD3, 0x99)
AMBER     = RGBColor(0xF7, 0xB9, 0x55)
RED       = RGBColor(0xF0, 0x6C, 0x6C)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)

FONT = "Arial"

prs = Presentation()
# 16:9
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.fill.solid(); bg.fill.fore_color.rgb = BG
    bg.line.fill.background()
    bg.shadow.inherit = False
    # send to back
    sp = bg._element
    sp.getparent().remove(sp)
    s.shapes._spTree.insert(2, sp)
    return s


def box(s, x, y, w, h, fill=None, line=None, line_w=1.0, radius=True):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shp = s.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    if radius:
        try:
            shp.adjustments[0] = 0.06
        except Exception:
            pass
    return shp


def text(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space_after=6, line_spacing=1.0):
    """runs: list of paragraphs; each paragraph is list of (txt, size, color, bold, italic)."""
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        p.line_spacing = line_spacing
        for (txt, size, color, bold, italic) in para:
            r = p.add_run(); r.text = txt
            r.font.size = Pt(size); r.font.color.rgb = color
            r.font.bold = bold; r.font.italic = italic
            r.font.name = FONT
    return tb


def kicker(s, x, y, label):
    dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y + 0.04), Inches(0.10), Inches(0.10))
    dot.fill.solid(); dot.fill.fore_color.rgb = ACCENT; dot.line.fill.background(); dot.shadow.inherit = False
    text(s, x + 0.22, y - 0.05, 8, 0.3, [[(label.upper(), 12, ACCENT, True, False)]])


def footer(s, num):
    text(s, 0.55, 7.05, 6, 0.3, [[("VANTAGE", 10, TEXT_DIM, True, False),
                                   ("  ·  Team Sentinels", 10, TEXT_FAINT, False, False)]])
    text(s, 11.2, 7.05, 1.6, 0.3, [[(f"{num:02d} / 06", 10, TEXT_FAINT, False, False)]],
         align=PP_ALIGN.RIGHT)


def bullets(s, x, y, w, h, items, size=14, gap=8, marker="▸", marker_color=ACCENT):
    """items: list of list-of-runs (txt,bold). Renders one wrapped bullet per item."""
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    for i, parts in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap); p.space_before = Pt(0); p.line_spacing = 1.05
        m = p.add_run(); m.text = f"{marker} "
        m.font.size = Pt(size); m.font.color.rgb = marker_color; m.font.bold = True; m.font.name = FONT
        for (txt, bold) in parts:
            r = p.add_run(); r.text = txt
            r.font.size = Pt(size); r.font.color.rgb = TEXT if bold else TEXT_DIM
            r.font.bold = bold; r.font.name = FONT
    return tb


# =================================================================
# SLIDE 1 — TITLE
# =================================================================
s = slide()
# logo mark
lm = box(s, 0.75, 1.35, 0.75, 0.75, fill=ACCENT)
text(s, 0.75, 1.35, 0.75, 0.75, [[("V", 32, WHITE, True, False)]],
     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
kicker(s, 1.75, 1.45, "Team Sentinels")
text(s, 1.97, 1.72, 9, 0.3, [[("Web Data Unlocked Hackathon · lablab.ai × Bright Data", 12, TEXT_FAINT, False, False)]])

text(s, 0.75, 2.45, 11, 1.4, [[("Vantage", 66, TEXT, True, False)]])
text(s, 0.78, 3.85, 10.2, 1.4, [[
    ("Time-travel for the open web. ", 21, TEXT, True, False),
    ("Type a company domain and watch its story — exec exits, layoffs, lawsuits, "
     "sentiment, pricing — rebuild itself on a draggable timeline. Every claim cited "
     "to its real source.", 21, TEXT_DIM, False, False)]], line_spacing=1.25)

pills = ["🏁 Track 2 · Finance & Market Intelligence", "🌐 Powered by Bright Data MCP",
         "🧠 Extraction by Claude", "⏱️ Ingest → render in ≤10s"]
px = 0.75
for pt in pills:
    w = 0.42 + 0.092 * len(pt)
    p = box(s, px, 5.65, w, 0.5, fill=BG_ELEV, line=BORDER)
    text(s, px, 5.65, w, 0.5, [[(pt, 12, TEXT_DIM, False, False)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    px += w + 0.18
footer(s, 1)

# =================================================================
# SLIDE 2 — PROBLEM / WHY / FOR WHO
# =================================================================
s = slide()
kicker(s, 0.6, 0.55, "The problem")
text(s, 0.55, 0.85, 12, 0.7, [[("Preliminary diligence is a 40-hour manual slog", 32, TEXT, True, False)]])
text(s, 0.58, 1.55, 12, 0.5, [[("Analysts hand-pull LinkedIn, Glassdoor, news, filings, pricing & lawsuits — "
                                "then it's stale by the time the memo ships.", 15, TEXT_DIM, False, False)]],
     line_spacing=1.2)

cards2 = [
    (0.6,  "40+ hrs", RED,   "per target, per analyst, on grunt research before any real analysis begins"),
    (4.82, "$$$",     RED,   "associate hours billed at PE / IB rates — multiplied across every deal in the pipeline"),
    (9.04, "Static",  AMBER, "snapshots only — you see today, never the trajectory that predicts tomorrow"),
]
for x, stat, col, desc in cards2:
    box(s, x, 2.45, 3.7, 1.55, fill=BG_CARD, line=BORDER)
    text(s, x + 0.28, 2.65, 3.2, 0.7, [[(stat, 34, col, True, False)]])
    text(s, x + 0.28, 3.32, 3.2, 0.6, [[(desc, 12, TEXT_DIM, False, False)]], line_spacing=1.12)

box(s, 0.6, 4.25, 12.13, 2.45, fill=BG_CARD, line=BORDER)
text(s, 0.9, 4.45, 11, 0.4, [[("👥  Who feels this pain", 17, TEXT, True, False)]])
whos = ["PE associates", "IB analysts", "Hedge-fund researchers", "Corp-dev / M&A teams", "VC scouts", "Credit & risk desks"]
wx = 0.9
for w_ in whos:
    bw = 0.4 + 0.105 * len(w_)
    box(s, wx, 5.0, bw, 0.46, fill=BG_ELEV, line=BORDER)
    text(s, wx, 5.0, bw, 0.46, [[(w_, 12, TEXT_DIM, False, False)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    wx += bw + 0.16
text(s, 0.9, 5.7, 11.5, 0.9, [[
    ("Vantage saves the time and the spend: ", 14.5, TEXT, True, False),
    ("it compresses that 40-hour pull into a ~10-second, fully-cited, time-aware briefing — "
     "so the human spends their hours on judgment, not on copy-paste.", 14.5, TEXT_DIM, False, False)]],
     line_spacing=1.25)
footer(s, 2)

# =================================================================
# SLIDE 3 — SOLUTION + DEMO PLACEHOLDER
# =================================================================
s = slide()
kicker(s, 0.6, 0.55, "The solution")
text(s, 0.55, 0.85, 12, 0.7, [[("A temporal knowledge graph of any company", 32, TEXT, True, False)]])

bullets(s, 0.65, 1.95, 5.2, 3.6, [
    [("Type a domain", True), (" → multi-source ingest fires across the open web", False)],
    [("Drag the time slider", True), (" (30 / 90 / 180 days) and the graph mutates live", False)],
    [("Hit play", True), (" → watch a company's collapse or breakout unfold week-by-week", False)],
    [("Every node & fact is cited", True), (" to the exact scraped article URL", False)],
    [("Auto-memo", True), (" with risk narrative + IC questions, export to PDF", False)],
], size=15, gap=13)
text(s, 0.65, 5.75, 5.2, 1.2, [[("Bitemporal model: every fact carries valid_from / valid_to / "
                                 "observed_at / source — that's what makes “what did this look like at "
                                 "time t?” possible.", 11.5, TEXT_FAINT, False, True)]], line_spacing=1.2)

# demo screenshot placeholder
ph = box(s, 6.3, 1.95, 6.43, 4.7, fill=BG_ELEV, line=BORDER, line_w=1.5)
text(s, 6.3, 3.6, 6.43, 1.5, [
    [("🖼️", 36, TEXT_DIM, False, False)],
    [("[ Demo screenshot goes here ]", 16, TEXT_DIM, True, False)],
    [("Dashboard with the time slider + animated knowledge graph", 12, TEXT_FAINT, False, False)],
    [("Replace this box with a PNG of /account/spirit.com", 12, TEXT_FAINT, False, False)],
], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=4)
footer(s, 3)

# =================================================================
# SLIDE 4 — WHY UNIQUE / NO COMPETITION
# =================================================================
s = slide()
kicker(s, 0.6, 0.55, "Why it's different")
text(s, 0.55, 0.85, 12.4, 0.7, [[("Nobody treats the open web as a time-aware database", 31, TEXT, True, False)]])
text(s, 0.58, 1.55, 12, 0.5, [[("Incumbents read internal data rooms & pre-timestamped filings. "
                                "We turn the messy public web into a queryable timeline.", 15, TEXT_DIM, False, True)]],
     line_spacing=1.2)

box(s, 0.6, 2.35, 5.95, 1.65, fill=BG_CARD, line=BORDER)
text(s, 0.88, 2.55, 5.5, 0.4, [[("🗂️  Hebbia · Rogo · AlphaSense", 16, TEXT, True, False)]])
text(s, 0.88, 3.05, 5.5, 0.9, [[("Operate on internal data rooms and already-timestamped SEC filings. "
                                 "They answer “what's in these documents?” — and the open web stays invisible.",
                                 13, TEXT_DIM, False, False)]], line_spacing=1.18)

box(s, 6.78, 2.35, 5.95, 1.65, fill=BG_CARD, line=ACCENT, line_w=1.5)
text(s, 7.06, 2.55, 5.5, 0.4, [[("⏳  Vantage", 16, TEXT, True, False)]])
text(s, 7.06, 3.05, 5.5, 0.9, [[("Structures the live public web into a bitemporal graph you can scrub "
                                 "through time. We answer ", 13, TEXT_DIM, False, False),
                                ("“how did this company change, and when?”", 13, TEXT, True, False)]],
     line_spacing=1.18)

box(s, 0.6, 4.25, 12.13, 2.45, fill=BG_CARD, line=BORDER)
text(s, 0.9, 4.45, 11, 0.4, [[("⏱️  The time dimension is the moat", 17, TEXT, True, False)]])
bullets(s, 0.9, 5.0, 11.5, 1.6, [
    [("Trajectory > snapshot", True), (" — two senior exits in 90 days is a distress signal a single snapshot can't see", False)],
    [("Lead/lag patterns", True), (" — a CMO departure often precedes a sentiment dip ~45 days later; we make that visible", False)],
    [("Replayable history", True), (" — drag back and watch Glassdoor stars climb, hiring expand, then the collapse animate forward", False)],
    [("Provenance over time", True), (" — every edge knows when it became true and where it came from", False)],
], size=13.5, gap=7)
footer(s, 4)

# =================================================================
# SLIDE 5 — BRIGHT DATA + EVERYTHING IT DOES
# =================================================================
s = slide()
kicker(s, 0.6, 0.55, "Built on Bright Data")
text(s, 0.55, 0.85, 12.4, 0.7, [[("Without Bright Data, this collapses to a static snapshot", 30, TEXT, True, False)]])

bd = [
    (0.6,  1.65, "🔌  MCP over SSE", "one interface",
     "A single orchestration layer across SERP, Web Scraper, Web Unlocker, Scraping Browser & Datasets — fan-out ingest from one client."),
    (6.78, 1.65, "🛡️  Web Unlocker + residential proxies", "",
     "Uninterrupted coverage of Glassdoor, LinkedIn & G2 where DIY scrapers die within days. This is what keeps the pipeline alive."),
    (0.6,  3.25, "🔎  SERP API", "",
     "News across NYT / FT / Business Journals, plus PACER litigation lookups — surfaced and placed on the timeline."),
    (6.78, 3.25, "📚  Datasets — historical backfill", "",
     "Pre-built LinkedIn / Glassdoor / Crunchbase history backfills 90+ days on Day 1 — no waiting to accumulate signal."),
]
for x, y, title, badge, desc in bd:
    box(s, x, y, 5.95, 1.45, fill=BG_CARD, line=BORDER)
    runs = [(title, 14.5, TEXT, True, False)]
    if badge:
        runs.append(("   " + badge, 11, ACCENT, False, False))
    text(s, x + 0.26, y + 0.16, 5.5, 0.4, [runs])
    text(s, x + 0.26, y + 0.62, 5.5, 0.75, [[(desc, 12, TEXT_DIM, False, False)]], line_spacing=1.15)

box(s, 0.6, 4.95, 12.13, 1.85, fill=BG_CARD, line=BORDER)
text(s, 0.9, 5.12, 11, 0.4, [[("🧩  Everything Vantage does with that data", 16, TEXT, True, False)]])
cols = [
    (0.9, [
        [("Multi-source ingest", True), (" across 6 source families", False)],
        [("Claude tool-use extraction", True), (" → entities + relations + events", False)],
        [("Bitemporal graph", True), (" (graphology) merge & resolve", False)],
    ]),
    (5.0, [
        [("Time slider + animated replay", True), (" over the graph", False)],
        [("Red-flag detection", True), (" across 6 risk dimensions", False)],
        [("Δ Delta report", True), (" — what changed between two dates", False)],
    ]),
    (9.1, [
        [("Cited Q&A", True), (" — answers link back to source spans", False)],
        [("Auto diligence memo", True), (" + PDF export", False)],
        [("Live usage metering", True), (" + 6h disk cache for speed", False)],
    ]),
]
for x, items in cols:
    bullets(s, x, 5.6, 3.9, 1.2, items, size=11.5, gap=5)
footer(s, 5)

# =================================================================
# SLIDE 6 — DEMO + LINKS
# =================================================================
s = slide()
kicker(s, 0.6, 0.55, "See it live")
text(s, 0.55, 0.85, 12, 0.7, [[("Try Vantage in 10 seconds", 32, TEXT, True, False)]])
text(s, 0.58, 1.55, 12, 0.5, [[("Three pre-cached targets load instantly through the real pipeline; "
                                "any other domain is fetched live via Bright Data.", 15, TEXT_DIM, False, False)]],
     line_spacing=1.2)

links = [
    ("🎥  DEMO VIDEO", "[ paste your demo video link here ]"),
    ("🚀  LIVE APP", "[ paste your deployed URL here ]"),
    ("💻  GITHUB REPO", "[ paste your repo link here ]"),
]
ly = 2.4
for lab, val in links:
    box(s, 0.6, ly, 5.95, 1.25, fill=BG_CARD, line=BORDER)
    text(s, 0.9, ly + 0.2, 5.4, 0.3, [[(lab, 11.5, TEXT_FAINT, True, False)]])
    text(s, 0.9, ly + 0.6, 5.4, 0.5, [[(val, 13, ACCENT, False, False)]])
    ly += 1.45

box(s, 6.78, 2.4, 5.95, 4.05, fill=BG_CARD, line=BORDER)
text(s, 7.06, 2.6, 5.5, 0.4, [[("✓  What the judges will see", 17, TEXT, True, False)]])
bullets(s, 7.06, 3.2, 5.5, 2.0, [
    [("Drag the slider — the graph mutates across 180 days", False)],
    [("Hit play — a company's story animates week by week", False)],
    [("Click any node — the exact scraped quote + source URL", False)],
    [("Open the memo — cited risk narrative, export to PDF", False)],
], size=14, gap=11, marker="✓", marker_color=GREEN)
text(s, 7.06, 5.75, 5.5, 0.5, [[("Team Sentinels", 14, TEXT, True, False),
                                ("  ·  Track 2  ·  Web Data Unlocked Hackathon", 14, TEXT_DIM, False, False)]])
footer(s, 6)

# ---- save ----
out = os.path.join(os.path.dirname(__file__), "Vantage-Sentinels-Submission.pptx")
prs.save(out)
print("Saved:", out)
